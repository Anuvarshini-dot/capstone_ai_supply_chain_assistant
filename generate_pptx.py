"""
Generate a detailed PowerPoint presentation for the AI Supply Chain Assistant capstone project.
Run with: py -3 generate_pptx.py
Output: AI_Supply_Chain_Assistant_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ─────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
MID_BG     = RGBColor(0x1A, 0x2E, 0x44)   # steel blue-grey
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
GREEN      = RGBColor(0x2D, 0xC6, 0x5C)   # success green
AMBER      = RGBColor(0xF7, 0xB7, 0x31)   # warning amber
RED        = RGBColor(0xEF, 0x47, 0x4A)   # danger red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xC4, 0xD8)
CARD_BG    = RGBColor(0x16, 0x28, 0x3A)   # card background

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


# ── Drawing helpers ───────────────────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor = None,
             line: RGBColor = None, line_width: float = 1.0) -> object:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.color.rgb = line if line else RGBColor(0, 0, 0)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.width = Pt(line_width)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color: RGBColor = WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_label_in_rect(slide, text, left, top, width, height,
                      bg: RGBColor = CARD_BG, font_size=11,
                      bold=False, font_color: RGBColor = WHITE,
                      line: RGBColor = ACCENT, line_width=1.5):
    add_rect(slide, left, top, width, height, fill=bg, line=line, line_width=line_width)
    add_text_box(slide, text, left, top, width, height,
                 font_size=font_size, bold=bold, color=font_color,
                 align=PP_ALIGN.CENTER)


def add_connector_arrow(slide, x1, y1, x2, y2, color: RGBColor = ACCENT, width=2.0):
    """Draw a straight connector line with an arrowhead."""
    from pptx.oxml.ns import qn
    from lxml import etree
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)

    # Top accent bar
    add_rect(slide, 0, 0, W, Inches(0.12), fill=ACCENT)

    # Side accent stripe
    add_rect(slide, 0, Inches(0.12), Inches(0.08), H - Inches(0.12), fill=ACCENT)

    # Main title
    add_text_box(slide, "AI-Powered Supply Chain\nIntelligence Assistant",
                 Inches(1.0), Inches(1.8), Inches(8.5), Inches(2.2),
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide, "Capstone Project — Architecture & Node Deep-Dive",
                 Inches(1.0), Inches(3.9), Inches(8.5), Inches(0.6),
                 font_size=20, bold=False, color=ACCENT2, align=PP_ALIGN.LEFT)

    # Bottom meta info
    add_text_box(slide,
                 "Built on  LangGraph  ·  FastAPI  ·  ChromaDB  ·  SQLite  ·  React",
                 Inches(1.0), Inches(5.0), Inches(10.0), Inches(0.5),
                 font_size=14, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    add_text_box(slide, "Presenter: Anuvarshini S S  |  FDE Capstone  |  2026",
                 Inches(1.0), Inches(5.6), Inches(10.0), Inches(0.4),
                 font_size=12, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # Decorative right graphic — stacked boxes
    for i, (label, col) in enumerate([
        ("Supplier Risk", RED), ("Shipment", AMBER), ("Inventory", GREEN),
    ]):
        add_label_in_rect(slide, label,
                          Inches(10.5), Inches(2.5 + i * 1.1),
                          Inches(2.3), Inches(0.75),
                          bg=col, font_size=13, bold=True,
                          font_color=WHITE, line=WHITE, line_width=0)

    # Bottom bar
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), fill=ACCENT)
    return slide


def slide_agenda(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), fill=ACCENT)

    add_text_box(slide, "Agenda", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=30, bold=True, color=WHITE)

    items = [
        ("01", "Project Overview & Problem Statement"),
        ("02", "System Architecture — End-to-End Flow"),
        ("03", "Node Deep-Dive: Input Guardrail & General Check"),
        ("04", "Node Deep-Dive: NL-to-SQL Agent"),
        ("05", "Node Deep-Dive: Specialist Classifier & Retrieval"),
        ("06", "Node Deep-Dive: Orchestrator & Specialist Agents"),
        ("07", "Node Deep-Dive: Summary & Recommendation Engine"),
        ("08", "Technology Stack & Data Model"),
        ("09", "Key Metrics & Business Value"),
        ("10", "Live Demo Walk-Through"),
    ]
    col1 = items[:5]
    col2 = items[5:]

    for i, (num, title) in enumerate(col1):
        add_text_box(slide, num, Inches(0.6), Inches(1.3 + i * 1.1),
                     Inches(0.55), Inches(0.6), font_size=22, bold=True, color=ACCENT)
        add_text_box(slide, title, Inches(1.25), Inches(1.3 + i * 1.1),
                     Inches(5.8), Inches(0.6), font_size=15, color=WHITE)

    for i, (num, title) in enumerate(col2):
        add_text_box(slide, num, Inches(7.3), Inches(1.3 + i * 1.1),
                     Inches(0.55), Inches(0.6), font_size=22, bold=True, color=ACCENT)
        add_text_box(slide, title, Inches(7.95), Inches(1.3 + i * 1.1),
                     Inches(5.0), Inches(0.6), font_size=15, color=WHITE)


def slide_overview(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), fill=ACCENT)

    add_text_box(slide, "Project Overview", Inches(0.5), Inches(0.25),
                 Inches(12), Inches(0.6), font_size=28, bold=True, color=WHITE)

    # Problem box
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(5.8), Inches(2.7),
             fill=CARD_BG, line=RED, line_width=1.5)
    add_text_box(slide, "The Problem", Inches(0.55), Inches(1.2),
                 Inches(5.5), Inches(0.45), font_size=16, bold=True, color=RED)
    problems = [
        "• Supply chain data is siloed across SQL databases, documents & logs",
        "• Analysts need hours to correlate supplier risk ↔ inventory ↔ shipment delays",
        "• No single view for proactive risk detection & recommendations",
        "• Business users can't query raw data without SQL knowledge",
    ]
    for i, p in enumerate(problems):
        add_text_box(slide, p, Inches(0.6), Inches(1.7 + i * 0.46),
                     Inches(5.5), Inches(0.45), font_size=12, color=LIGHT_GREY)

    # Solution box
    add_rect(slide, Inches(6.7), Inches(1.1), Inches(6.2), Inches(2.7),
             fill=CARD_BG, line=GREEN, line_width=1.5)
    add_text_box(slide, "Our Solution", Inches(6.85), Inches(1.2),
                 Inches(5.9), Inches(0.45), font_size=16, bold=True, color=GREEN)
    solutions = [
        "• Natural language interface — no SQL needed",
        "• Multi-agent pipeline: SQL + vector search + specialist AI agents",
        "• Real-time anomaly correlation across supplier / shipment / inventory",
        "• Automated risk detection + actionable recommendations",
    ]
    for i, s in enumerate(solutions):
        add_text_box(slide, s, Inches(6.9), Inches(1.7 + i * 0.46),
                     Inches(5.9), Inches(0.45), font_size=12, color=LIGHT_GREY)

    # Impact metrics
    metrics = [
        ("3 Databases", "SQLite + ChromaDB\n+ BM25 Index"),
        ("6 AI Agents", "Supplier · Shipment\nInventory · NL-SQL\nSummary · Recommend"),
        ("9 Graph Nodes", "Full LangGraph\nstate machine"),
        ("Hybrid RAG", "Vector + BM25\n+ Reranker"),
    ]
    for i, (title, sub) in enumerate(metrics):
        x = Inches(0.4 + i * 3.2)
        add_rect(slide, x, Inches(4.15), Inches(3.0), Inches(2.8),
                 fill=MID_BG, line=ACCENT, line_width=1.2)
        add_text_box(slide, title, x, Inches(4.35), Inches(3.0), Inches(0.6),
                     font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide, sub, x, Inches(4.95), Inches(3.0), Inches(1.6),
                     font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Main end-to-end flow architecture slide."""
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), fill=ACCENT)

    add_text_box(slide, "System Architecture — End-to-End Query Flow",
                 Inches(0.4), Inches(0.22), Inches(12), Inches(0.55),
                 font_size=26, bold=True, color=WHITE)

    # ── Lane backgrounds ──────────────────────────────────────────────────────
    # User lane
    add_rect(slide, Inches(0.15), Inches(0.95), Inches(12.9), Inches(1.05),
             fill=RGBColor(0x10, 0x22, 0x35), line=ACCENT2, line_width=0.5)
    add_text_box(slide, "USER", Inches(0.22), Inches(0.97), Inches(1.2), Inches(0.4),
                 font_size=9, bold=True, color=ACCENT2)

    # Backend pipeline lane
    add_rect(slide, Inches(0.15), Inches(2.1), Inches(12.9), Inches(3.2),
             fill=RGBColor(0x0F, 0x20, 0x30), line=ACCENT2, line_width=0.5)
    add_text_box(slide, "LANGGRAPH PIPELINE", Inches(0.22), Inches(2.12),
                 Inches(2.0), Inches(0.4), font_size=9, bold=True, color=ACCENT2)

    # Data layer
    add_rect(slide, Inches(0.15), Inches(5.45), Inches(12.9), Inches(1.85),
             fill=RGBColor(0x0C, 0x1A, 0x28), line=ACCENT2, line_width=0.5)
    add_text_box(slide, "DATA LAYER", Inches(0.22), Inches(5.47),
                 Inches(1.5), Inches(0.4), font_size=9, bold=True, color=ACCENT2)

    # ── User layer ────────────────────────────────────────────────────────────
    add_label_in_rect(slide, "React\nFrontend",
                      Inches(0.4), Inches(1.05), Inches(1.5), Inches(0.75),
                      bg=MID_BG, font_size=11, bold=True, line=ACCENT)
    add_label_in_rect(slide, "FastAPI\n/query",
                      Inches(2.2), Inches(1.05), Inches(1.5), Inches(0.75),
                      bg=MID_BG, font_size=11, bold=True, line=ACCENT)

    add_connector_arrow(slide, Inches(1.9), Inches(1.42), Inches(2.2), Inches(1.42))

    # Annotation
    add_text_box(slide, "POST /api/query  {query, filters, top_k}",
                 Inches(1.93), Inches(0.97), Inches(3.0), Inches(0.35),
                 font_size=8, color=LIGHT_GREY, italic=True)

    # ── Pipeline nodes ────────────────────────────────────────────────────────
    nodes = [
        ("🛡️ Input\nGuardrail",   CARD_BG, AMBER),
        ("🔍 General\nCheck",     CARD_BG, ACCENT),
        ("🗃️ NL→SQL\nAgent",      CARD_BG, GREEN),
        ("🎯 Specialist\nClassify", CARD_BG, ACCENT),
        ("🗄️ Vector\nRetrieval",  CARD_BG, ACCENT),
        ("🎛️ Orchestrator\n+ Agents", CARD_BG, RGBColor(0xAA, 0x77, 0xFF)),
        ("✦ Summary\nNode",       CARD_BG, ACCENT),
        ("💡 Recommend\nEngine",  CARD_BG, GREEN),
    ]

    node_w = Inches(1.45)
    node_h = Inches(0.88)
    gap    = Inches(0.17)
    start_x = Inches(0.32)
    node_y  = Inches(2.45)

    node_xs = []
    for i, (label, bg, line_col) in enumerate(nodes):
        x = start_x + i * (node_w + gap)
        node_xs.append(x + node_w / 2)
        add_label_in_rect(slide, label, x, node_y, node_w, node_h,
                          bg=bg, font_size=10, bold=True,
                          font_color=WHITE, line=line_col, line_width=1.8)

    # Arrows between nodes
    arrow_y = node_y + node_h / 2
    for i in range(len(nodes) - 1):
        x1 = start_x + i * (node_w + gap) + node_w
        x2 = start_x + (i + 1) * (node_w + gap)
        add_connector_arrow(slide, x1, arrow_y, x2, arrow_y, color=ACCENT, width=1.5)

    # Conditional branch labels
    # Guardrail → END
    add_text_box(slide, "invalid →\nEND", Inches(0.55), Inches(3.45),
                 Inches(1.1), Inches(0.55), font_size=8, color=RED, italic=True)
    add_connector_arrow(slide, Inches(1.05), Inches(3.33), Inches(1.05), Inches(3.8),
                        color=RED, width=1.0)

    # General check → general node branch
    add_text_box(slide, "off-topic →\nGeneral QA",
                 Inches(2.05), Inches(3.45), Inches(1.5), Inches(0.5),
                 font_size=8, color=AMBER, italic=True)

    # Classify → recommend shortcut
    add_text_box(slide, "no agents →\nskip to →",
                 Inches(5.9), Inches(3.45), Inches(1.5), Inches(0.5),
                 font_size=8, color=AMBER, italic=True)

    # ── Specialist agents sub-row ─────────────────────────────────────────────
    agents = [
        ("🏭 Supplier\nRisk Agent",   RED),
        ("🚢 Shipment\nAgent",        AMBER),
        ("📦 Inventory\nAgent",       GREEN),
    ]
    agt_w   = Inches(1.65)
    agt_h   = Inches(0.75)
    agt_y   = Inches(3.7)
    agt_xs_start = Inches(7.25)
    for i, (lbl, col) in enumerate(agents):
        ax = agt_xs_start + i * (agt_w + Inches(0.12))
        add_label_in_rect(slide, lbl, ax, agt_y, agt_w, agt_h,
                          bg=MID_BG, font_size=10, bold=True,
                          font_color=WHITE, line=col, line_width=2.0)
        add_connector_arrow(slide,
                            agt_xs_start + Inches(2.7), node_y + node_h,
                            ax + agt_w / 2, agt_y, color=col, width=1.2)

    # ── Data layer ────────────────────────────────────────────────────────────
    data_stores = [
        ("SQLite DB\n(Structured)", GREEN,   Inches(0.5)),
        ("ChromaDB\n(Vector Store)", ACCENT,  Inches(3.4)),
        ("BM25 Index\n(Keyword)", ACCENT2,  Inches(6.3)),
        ("LLM Gateway\n(OpenAI)", AMBER,   Inches(9.2)),
    ]
    for label, col, dx in data_stores:
        add_label_in_rect(slide, label, dx, Inches(5.65), Inches(2.5), Inches(0.9),
                          bg=CARD_BG, font_size=11, bold=True,
                          font_color=WHITE, line=col, line_width=1.5)
        # arrow from pipeline to data layer
        add_connector_arrow(slide, dx + Inches(1.25), Inches(5.45),
                            dx + Inches(1.25), Inches(5.65),
                            color=col, width=1.0)


def slide_node_detail(prs, icon, title, subtitle, role_text,
                      inputs, outputs, calls, notes, accent_col: RGBColor = ACCENT):
    """Reusable per-node detail slide."""
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=accent_col)

    # Header stripe
    add_rect(slide, 0, Inches(0.1), W, Inches(1.05), fill=MID_BG)
    add_text_box(slide, f"{icon}  {title}",
                 Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.65),
                 font_size=26, bold=True, color=WHITE)
    add_text_box(slide, subtitle,
                 Inches(0.4), Inches(0.78), Inches(9.0), Inches(0.4),
                 font_size=13, color=ACCENT2, italic=True)

    # Role summary
    add_rect(slide, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.72),
             fill=CARD_BG, line=accent_col, line_width=1.2)
    add_text_box(slide, role_text,
                 Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.65),
                 font_size=12, color=WHITE)

    # Three columns: Inputs | Calls | Outputs
    col_w = Inches(4.1)
    col_tops = [Inches(2.25), Inches(2.25), Inches(2.25)]
    col_xs   = [Inches(0.3), Inches(4.6), Inches(8.9)]

    headers_info = [
        ("INPUTS", inputs,  ACCENT2),
        ("CALLS / INVOKES", calls, AMBER),
        ("OUTPUTS", outputs, GREEN),
    ]
    for (hdr, items, hdr_col), cx, ct in zip(headers_info, col_xs, col_tops):
        add_rect(slide, cx, ct, col_w, Inches(0.38), fill=CARD_BG, line=hdr_col, line_width=1.5)
        add_text_box(slide, hdr, cx, ct + Inches(0.05), col_w, Inches(0.3),
                     font_size=11, bold=True, color=hdr_col, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            item_y = ct + Inches(0.45) + j * Inches(0.42)
            add_rect(slide, cx + Inches(0.05), item_y,
                     col_w - Inches(0.1), Inches(0.38),
                     fill=RGBColor(0x12, 0x22, 0x32), line=None)
            add_text_box(slide, f"  {item}", cx + Inches(0.1), item_y,
                         col_w - Inches(0.2), Inches(0.38),
                         font_size=10.5, color=LIGHT_GREY)

    # Notes / Decision logic
    if notes:
        note_y = Inches(5.5)
        add_rect(slide, Inches(0.3), note_y, Inches(12.7), Inches(1.7),
                 fill=CARD_BG, line=accent_col, line_width=1.0)
        add_text_box(slide, "Decision Logic / Notes",
                     Inches(0.45), note_y + Inches(0.05), Inches(4.0), Inches(0.35),
                     font_size=11, bold=True, color=accent_col)
        for k, note in enumerate(notes):
            add_text_box(slide, f"• {note}",
                         Inches(0.5), note_y + Inches(0.42) + k * Inches(0.38),
                         Inches(12.2), Inches(0.36), font_size=11, color=LIGHT_GREY)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=accent_col)


def slide_orchestrator(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=RGBColor(0xAA, 0x77, 0xFF))

    add_rect(slide, 0, Inches(0.1), W, Inches(1.0), fill=MID_BG)
    add_text_box(slide, "🎛️  Orchestrator Node — Specialist Agent Pipeline",
                 Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.6),
                 font_size=24, bold=True, color=WHITE)
    add_text_box(slide,
                 "Runs specialist agents in classifier-determined order; each agent receives focused sub-queries + targeted documents",
                 Inches(0.4), Inches(0.75), Inches(12.0), Inches(0.38),
                 font_size=12, color=ACCENT2, italic=True)

    # ── Execution flow ────────────────────────────────────────────────────────
    flow_items = [
        ("1. Read ordered agent list from classify_node", ACCENT),
        ("2. For each agent: call _targeted_docs() → fetch profile + shipment docs via hybrid_search + reranker", ACCENT),
        ("3. Call _run_agent() with: focused sub_query, targeted docs, prior agent findings (context chaining)", ACCENT),
        ("4. Store result in agent_findings[agent_name]", ACCENT),
        ("5. Append rich log entry: risk_level, confidence, affected_suppliers / warehouses_at_risk", ACCENT),
    ]
    add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.8),
             fill=CARD_BG, line=RGBColor(0xAA, 0x77, 0xFF), line_width=1.5)
    add_text_box(slide, "Execution Sequence",
                 Inches(0.45), Inches(1.35), Inches(6.0), Inches(0.4),
                 font_size=13, bold=True, color=RGBColor(0xAA, 0x77, 0xFF))
    for i, (txt, col) in enumerate(flow_items):
        add_text_box(slide, txt, Inches(0.5), Inches(1.82) + Inches(0.58) * i,
                     Inches(5.9), Inches(0.52), font_size=11, color=LIGHT_GREY)

    # ── Agents ────────────────────────────────────────────────────────────────
    agents = [
        ("🏭 Supplier Risk Agent", RED,
         ["Supplier profiles: delay rate, defect rate, reliability",
          "Identifies high-risk / underperforming suppliers",
          "Outputs: risk_level, confidence, affected_suppliers"]),
        ("🚢 Shipment Agent", AMBER,
         ["Shipment events: delays, carriers, routes, regions",
          "Detects at-risk or delayed shipments by mode/carrier",
          "Outputs: risk_level, confidence, top findings"]),
        ("📦 Inventory Agent", GREEN,
         ["Warehouse profiles: capacity, utilisation, region",
          "Identifies stockout risk, low inventory, reorder needs",
          "Outputs: risk_level, confidence, warehouses_at_risk"]),
    ]
    for i, (lbl, col, bullets) in enumerate(agents):
        ax = Inches(6.8) + i * Inches(2.18)
        add_rect(slide, ax, Inches(1.3), Inches(2.05), Inches(3.85),
                 fill=CARD_BG, line=col, line_width=2.0)
        add_text_box(slide, lbl, ax, Inches(1.35), Inches(2.05), Inches(0.52),
                     font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            add_text_box(slide, f"• {b}", ax + Inches(0.1),
                         Inches(1.92) + j * Inches(0.75), Inches(1.88), Inches(0.72),
                         font_size=9.5, color=LIGHT_GREY)

    # Context chaining note
    add_rect(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.95),
             fill=RGBColor(0x10, 0x1E, 0x2D), line=AMBER, line_width=1.2)
    add_text_box(slide, "Context Chaining (Key Design Pattern)",
                 Inches(0.5), Inches(5.35), Inches(6.0), Inches(0.38),
                 font_size=12, bold=True, color=AMBER)
    add_text_box(slide,
                 "Each agent receives findings_so_far (all prior agents' outputs) as prior_findings input.\n"
                 "Example: Inventory Agent sees Supplier Risk findings → can correlate supplier delays with stockout risk.\n"
                 "This enables cross-domain anomaly detection without re-querying the same data.",
                 Inches(0.5), Inches(5.75), Inches(12.2), Inches(1.4),
                 font_size=11, color=LIGHT_GREY)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1),
             fill=RGBColor(0xAA, 0x77, 0xFF))


def slide_retrieval_deep(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=ACCENT)
    add_rect(slide, 0, Inches(0.1), W, Inches(1.0), fill=MID_BG)

    add_text_box(slide, "🗄️  Vector Retrieval Node — Hybrid RAG Pipeline",
                 Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.6),
                 font_size=24, bold=True, color=WHITE)
    add_text_box(slide, "3-step retrieval: profile docs (aggregated) + shipment events + reranker fusion",
                 Inches(0.4), Inches(0.75), Inches(12.0), Inches(0.38),
                 font_size=12, color=ACCENT2, italic=True)

    steps = [
        ("Step 1 — Profile Docs", ACCENT,
         "hybrid_search(query, filters={doc_type: supplier_profile | warehouse_profile, entity_names})\n"
         "Fetches pre-aggregated performance profiles for SQL-identified entities.\n"
         "Contains: delay rates, avg delay, inventory impact, reliability scores."),
        ("Step 2 — Shipment Events", GREEN,
         "hybrid_search(query, top_k=15, entity_filter)\n"
         "Fetches individual shipment/incident records for granular context.\n"
         "Fallback: relaxes entity filter if < 2 docs returned."),
        ("Step 3 — Reranker Fusion", AMBER,
         "rerank(query, shipment_docs, top_k=remaining)\n"
         "Cross-encoder reranker scores remaining slots after profiles guaranteed.\n"
         "Ensures most relevant docs fill remaining top_k budget."),
    ]
    for i, (title, col, desc) in enumerate(steps):
        y = Inches(1.35) + i * Inches(1.78)
        add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(1.6),
                 fill=CARD_BG, line=col, line_width=1.5)
        add_text_box(slide, title, Inches(0.5), y + Inches(0.08),
                     Inches(5.0), Inches(0.42), font_size=14, bold=True, color=col)
        add_text_box(slide, desc, Inches(0.5), y + Inches(0.52),
                     Inches(12.0), Inches(1.0), font_size=11, color=LIGHT_GREY)

    # hybrid_search explanation
    add_rect(slide, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
             fill=RGBColor(0x10, 0x1E, 0x2D), line=ACCENT2, line_width=1.0)
    add_text_box(slide,
                 "hybrid_search = ChromaDB vector similarity (embeddings) + BM25 keyword index, "
                 "merged with Reciprocal Rank Fusion (RRF). Captures both semantic and exact-keyword relevance.",
                 Inches(0.5), Inches(6.78), Inches(12.2), Inches(0.55),
                 font_size=10.5, color=LIGHT_GREY, italic=True)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=ACCENT)


def slide_tech_stack(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=ACCENT)
    add_rect(slide, 0, Inches(0.1), W, Inches(0.9), fill=MID_BG)

    add_text_box(slide, "Technology Stack & Data Architecture",
                 Inches(0.4), Inches(0.18), Inches(12.0), Inches(0.6),
                 font_size=26, bold=True, color=WHITE)

    categories = [
        ("Frontend", ACCENT, [
            "React 18 + Vite",
            "React Router",
            "Custom Hooks (useSupplyChain)",
            "Axios REST client",
            "CSS Modules",
        ]),
        ("Backend API", GREEN, [
            "FastAPI (Python 3.12)",
            "Pydantic schemas",
            "Async endpoints",
            "LangGraph StateGraph",
            "Evaluation runner",
        ]),
        ("AI / LLM", RGBColor(0xAA, 0x77, 0xFF), [
            "OpenAI via internal gateway",
            "6 specialist agents",
            "NL-to-SQL chat loop",
            "Hybrid RAG pipeline",
            "Cross-encoder reranker",
        ]),
        ("Data Layer", AMBER, [
            "SQLite (structured ops data)",
            "ChromaDB (vector store)",
            "BM25 keyword index",
            "Supplier / Warehouse profiles",
            "Shipment event documents",
        ]),
    ]

    cat_w = Inches(3.1)
    for i, (cat, col, items) in enumerate(categories):
        cx = Inches(0.25) + i * (cat_w + Inches(0.15))
        add_rect(slide, cx, Inches(1.2), cat_w, Inches(5.85),
                 fill=CARD_BG, line=col, line_width=2.0)
        add_text_box(slide, cat, cx, Inches(1.26), cat_w, Inches(0.45),
                     font_size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, f"• {item}",
                         cx + Inches(0.15), Inches(1.75) + j * Inches(0.95),
                         cat_w - Inches(0.25), Inches(0.85),
                         font_size=12, color=LIGHT_GREY)

    # DB schema summary
    add_rect(slide, Inches(0.25), Inches(7.12), Inches(12.65), Inches(0.25),
             fill=CARD_BG, line=ACCENT, line_width=0.8)
    add_text_box(slide,
                 "SQLite tables: suppliers · products · warehouses · inventory · shipments",
                 Inches(0.4), Inches(7.12), Inches(12.0), Inches(0.25),
                 font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=ACCENT)


def slide_anomaly_and_eval(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=GREEN)
    add_rect(slide, 0, Inches(0.1), W, Inches(0.95), fill=MID_BG)

    add_text_box(slide, "Anomaly Correlation & Evaluation Pipeline",
                 Inches(0.4), Inches(0.18), Inches(12.0), Inches(0.58),
                 font_size=26, bold=True, color=WHITE)
    add_text_box(slide,
                 "Cross-agent anomaly detection + automated RAG faithfulness & relevance scoring",
                 Inches(0.4), Inches(0.72), Inches(12.0), Inches(0.36),
                 font_size=12, color=ACCENT2, italic=True)

    anomalies = [
        ("Supplier → Inventory Cascade", RED,
         "Triggered when: supplier risk = HIGH and inventory risk = MEDIUM|HIGH\n"
         "Meaning: Supplier delivery delays are correlating with warehouse depletion.\n"
         "Severity: HIGH — requires immediate dual remediation."),
        ("Shipment + Supplier Compound", AMBER,
         "Triggered when: shipment risk = HIGH and supplier risk = MEDIUM|HIGH\n"
         "Meaning: Concurrent disruptions create a compounded bottleneck across two domains.\n"
         "Severity: HIGH — single-point mitigation is insufficient."),
        ("Inventory + Shipment Stockout Risk", GREEN,
         "Triggered when: inventory risk = HIGH and shipment risk = MEDIUM|HIGH\n"
         "Meaning: Delayed inbound shipments compounding already-low stock levels.\n"
         "Severity: MEDIUM — stockout window is narrowing."),
    ]
    for i, (title, col, desc) in enumerate(anomalies):
        y = Inches(1.28) + i * Inches(1.55)
        add_rect(slide, Inches(0.3), y, Inches(12.65), Inches(1.38),
                 fill=CARD_BG, line=col, line_width=1.8)
        add_text_box(slide, title, Inches(0.5), y + Inches(0.07),
                     Inches(6.0), Inches(0.38), font_size=13, bold=True, color=col)
        add_text_box(slide, desc, Inches(0.5), y + Inches(0.48),
                     Inches(12.0), Inches(0.85), font_size=11, color=LIGHT_GREY)

    # Evaluation metrics
    add_rect(slide, Inches(0.3), Inches(5.98), Inches(12.65), Inches(1.35),
             fill=CARD_BG, line=ACCENT2, line_width=1.2)
    add_text_box(slide, "Automated Evaluation (evaluation/runner.py)",
                 Inches(0.5), Inches(6.02), Inches(6.0), Inches(0.38),
                 font_size=13, bold=True, color=ACCENT2)
    evals = [
        "Faithfulness — answer grounded in retrieved docs/SQL data",
        "Context Relevance — retrieved docs match the query",
        "Answer Confidence — aggregated agent confidence scores",
        "Anomaly Count — cross-domain correlation detections",
    ]
    for k, e in enumerate(evals):
        cx = Inches(0.5) if k < 2 else Inches(6.8)
        cy = Inches(6.45) + (k % 2) * Inches(0.4)
        add_text_box(slide, f"• {e}", cx, cy, Inches(6.0), Inches(0.38),
                     font_size=11, color=LIGHT_GREY)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=GREEN)


def slide_business_value(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=ACCENT)
    add_rect(slide, 0, Inches(0.1), W, Inches(0.9), fill=MID_BG)

    add_text_box(slide, "Business Value & Key Capabilities",
                 Inches(0.4), Inches(0.18), Inches(12.0), Inches(0.58),
                 font_size=26, bold=True, color=WHITE)

    cards = [
        ("Natural Language Interface", ACCENT,
         "Business users query supply chain data in plain English — no SQL expertise required.\n"
         "Input guardrail ensures security: PII blocking, length & format validation."),
        ("Real-Time Risk Intelligence", RED,
         "Three specialist AI agents run in parallel-aware, context-chained order.\n"
         "Risk levels (high/medium/low) surfaced per domain with confidence scores."),
        ("Cross-Domain Anomaly Detection", AMBER,
         "3 anomaly correlation patterns detect cascade failures across supplier, shipment & inventory domains.\n"
         "Prevents siloed analysis that misses multi-domain compounding effects."),
        ("Actionable Recommendations", GREEN,
         "Recommendation Engine generates domain-specific, prioritised actions.\n"
         "Based on all agent findings — not just a single data source."),
        ("Full Observability", RGBColor(0xAA, 0x77, 0xFF),
         "Every pipeline step logged with: timing (ms), confidence, risk level, top finding.\n"
         "Frontend Agent Trace panel lets users inspect every decision the system made."),
        ("Hybrid RAG — Best of Both Worlds", ACCENT2,
         "ChromaDB vector search (semantic) + BM25 (keyword) + cross-encoder reranker.\n"
         "Handles both vague semantic queries and precise entity-specific lookups."),
    ]

    card_w = Inches(4.15)
    card_h = Inches(2.25)
    for i, (title, col, desc) in enumerate(cards):
        cx = Inches(0.28) + (i % 3) * (card_w + Inches(0.18))
        cy = Inches(1.18) + (i // 3) * (card_h + Inches(0.18))
        add_rect(slide, cx, cy, card_w, card_h, fill=CARD_BG, line=col, line_width=2.0)
        add_text_box(slide, title, cx + Inches(0.15), cy + Inches(0.1),
                     card_w - Inches(0.2), Inches(0.45),
                     font_size=13, bold=True, color=col)
        add_text_box(slide, desc, cx + Inches(0.15), cy + Inches(0.58),
                     card_w - Inches(0.25), Inches(1.55),
                     font_size=10.5, color=LIGHT_GREY)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=ACCENT)


def slide_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=ACCENT)

    add_text_box(slide, "Live Demo Walk-Through",
                 Inches(0.5), Inches(1.2), Inches(12.0), Inches(0.8),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sample_queries = [
        ("Supplier Risk", RED,
         "\"Which suppliers have the highest defect rates and are causing inventory issues?\"",
         "Triggers: NL→SQL → Supplier Risk Agent → Inventory Agent → Anomaly Correlation"),
        ("Shipment Delay", AMBER,
         "\"Show me all shipments delayed by more than 5 days in the LATAM region\"",
         "Triggers: NL→SQL → Shipment Agent → Risk Assessment → Recommendations"),
        ("Inventory Health", GREEN,
         "\"Which warehouses are at critical stock levels and which suppliers serve them?\"",
         "Triggers: NL→SQL → Inventory Agent → Supplier Agent → Cascade Anomaly Check"),
    ]
    for i, (label, col, query, trace) in enumerate(sample_queries):
        y = Inches(2.35) + i * Inches(1.5)
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(1.3),
                 fill=CARD_BG, line=col, line_width=1.8)
        add_text_box(slide, label, Inches(0.6), y + Inches(0.07),
                     Inches(2.5), Inches(0.38), font_size=12, bold=True, color=col)
        add_text_box(slide, query, Inches(0.6), y + Inches(0.48),
                     Inches(12.0), Inches(0.4), font_size=12, color=WHITE, italic=True)
        add_text_box(slide, trace, Inches(0.6), y + Inches(0.88),
                     Inches(12.0), Inches(0.35), font_size=10, color=LIGHT_GREY)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=ACCENT)


def slide_closing(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.1), fill=ACCENT)

    add_text_box(slide, "Thank You",
                 Inches(1.0), Inches(1.8), Inches(10.0), Inches(1.2),
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "AI-Powered Supply Chain Intelligence Assistant",
                 Inches(1.0), Inches(3.1), Inches(10.0), Inches(0.55),
                 font_size=20, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Questions & Discussion",
                 Inches(1.0), Inches(3.75), Inches(10.0), Inches(0.5),
                 font_size=16, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    add_text_box(slide, "Anuvarshini S S  |  FDE Capstone  |  2026",
                 Inches(1.0), Inches(5.2), Inches(10.0), Inches(0.4),
                 font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=ACCENT)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)
    slide_overview(prs)
    slide_architecture(prs)

    # Per-node slides
    slide_node_detail(
        prs,
        icon="🛡️", title="Input Guardrail Node",
        subtitle="First gate — validates every query before any AI or database work begins",
        accent_col=AMBER,
        role_text=(
            "Calls validate_query() from guardrails/input_validator.py. "
            "Checks: query length (min/max), PII patterns (emails, phone numbers, SSNs), "
            "and disallowed content. If validation fails, sets answer = error message and short-circuits the entire graph."
        ),
        inputs=[
            "query (str) — raw user input",
            "execution_log (list) — accumulated log",
        ],
        outputs=[
            "validation_passed: True | False",
            "answer (str) — error message if rejected",
            "execution_log — updated with guardrail result",
        ],
        calls=[
            "validate_query(query) → raises ValidationError",
            "route_after_guardrail() → 'end' | 'general_check_node'",
        ],
        notes=[
            "If validation_passed = False → graph routes to END immediately; no LLM or DB call is made",
            "Execution log always gets an entry regardless of pass/fail — full observability",
            "PII blocking prevents sensitive data from reaching external LLM gateway",
        ],
    )

    slide_node_detail(
        prs,
        icon="🔍", title="General Check Node",
        subtitle="Binary router — supply chain topic or off-topic? Decides the pipeline path",
        accent_col=ACCENT,
        role_text=(
            "Calls the LLM with a strict system prompt to return {\"is_supply_chain\": true/false}. "
            "Has a keyword safety net (_SC_KEYWORDS set with 30+ supply chain terms) that overrides "
            "an LLM 'false' if any keyword matches — prevents false negatives on legitimate queries."
        ),
        inputs=[
            "query (str)",
            "execution_log",
        ],
        outputs=[
            "routed_agents: ['supply_chain'] | ['general']",
            "execution_log — updated",
        ],
        calls=[
            "chat([system, user]) → LLM classification",
            "json.loads() → parse {is_supply_chain: bool}",
            "keyword override: any(kw in query for kw in _SC_KEYWORDS)",
            "route_after_general_check() → 'general_node' | 'nlsql_node'",
        ],
        notes=[
            "Supply chain path → nlsql_node (structured data query always runs first)",
            "Off-topic path → general_node → BaseAgent.answer_general() → END",
            "Keyword safety net ensures supply chain queries with unusual phrasing are not lost",
            "Defaults is_supply_chain = True on any JSON parse error (fail-safe)",
        ],
    )

    slide_node_detail(
        prs,
        icon="🗃️", title="NL-to-SQL Agent Node",
        subtitle="Converts natural language query to SQL, executes against SQLite, extracts entities",
        accent_col=GREEN,
        role_text=(
            "NLSQLAgent runs a chat loop (up to 5 iterations): ask LLM for SQL → extract SQL block → "
            "execute against SQLite → feed result back to LLM → repeat until no more SQL blocks. "
            "Extracts named entities (supplier_names, warehouse_names, product_names) from results "
            "to enable targeted ChromaDB retrieval in subsequent nodes."
        ),
        inputs=[
            "query (str)",
        ],
        outputs=[
            "sql_result — SQL queries used",
            "sql_data — raw tabular output",
            "answer — LLM-generated factual answer",
            "sql_entities — {supplier_names, warehouse_names, product_names}",
            "agent_findings['nlsql'] — summary, confidence, findings",
            "confidence_score — 0.9 (high, structured data)",
        ],
        calls=[
            "chat(messages) — up to 5 LLM calls in loop",
            "sqlite3.connect(SQLITE_DB_PATH) → pd.read_sql_query(sql)",
            "_extract_sql(response) → regex for ```sql blocks",
            "_extract_entities(dfs) → pull entity names from results",
            "_build_display_facts(dfs) → top 3 rows × 2 cols for UI cards",
        ],
        notes=[
            "Schema injected into system prompt — LLM knows all 5 tables and join rules",
            "Multi-query loop handles complex questions requiring multiple JOINs",
            "sql_entities feeds directly into retrieve_node and orchestrator for targeted retrieval",
            "Confidence fixed at 0.9 — structured SQL data is highly reliable",
        ],
    )

    slide_node_detail(
        prs,
        icon="🎯", title="Specialist Classifier Node",
        subtitle="Decides which specialist agents run, in what order, with what focused sub-question",
        accent_col=ACCENT,
        role_text=(
            "LLM chooses 0–3 specialists (supplier, shipment, inventory) and writes a focused sub_query "
            "for each. The ordering rule: most data-rich agent for the query runs first. "
            "SQL entity context is injected so the LLM knows exact entity names and can write "
            "precise sub-queries (e.g. 'Los Angeles Fulfillment Hub' not just 'the warehouse')."
        ),
        inputs=[
            "query (str)",
            "sql_entities — entity names from NL→SQL",
            "execution_log",
        ],
        outputs=[
            "routed_agents — ordered list e.g. ['inventory', 'supplier']",
            "agent_sub_queries — {agent: focused_question}",
            "execution_log — updated with plan",
        ],
        calls=[
            "chat([system, user+entity_context]) → LLM classification",
            "json.loads() → parse {agents: [{name, sub_query}]}",
            "keyword fallback — if JSON parse fails, use domain keywords",
            "route_after_classify() → 'retrieve_node' | 'recommendation_node'",
        ],
        notes=[
            "If 0 agents selected → skip to recommendation_node (SQL answer is sufficient)",
            "Entity context injection: 'SQL identified suppliers: GlobalTech, Apex Mfg'",
            "Ordering matters — earlier agents' findings are fed to later agents as prior_findings",
            "Keyword fallback covers malformed LLM JSON responses gracefully",
        ],
    )

    slide_retrieval_deep(prs)
    slide_orchestrator(prs)

    slide_node_detail(
        prs,
        icon="✦", title="Summary Node",
        subtitle="Synthesises all agent findings into a coherent, human-readable narrative",
        accent_col=ACCENT2,
        role_text=(
            "SummaryAgent.summarize() receives the query and all agent_findings (nlsql, supplier, "
            "shipment, inventory) and produces a unified answer. Also runs _detect_anomaly_correlations() "
            "to identify cross-domain cascade patterns. Computes overall confidence as the mean of "
            "all individual agent confidence scores."
        ),
        inputs=[
            "query (str)",
            "agent_findings — all agent outputs keyed by name",
            "execution_log",
        ],
        outputs=[
            "answer — unified narrative summary",
            "anomaly_correlations — list of detected cascade patterns",
            "confidence_score — mean of all agent confidences",
            "execution_log — updated",
        ],
        calls=[
            "SummaryAgent().summarize(query, findings) → LLM synthesis",
            "_detect_anomaly_correlations(findings) → cross-agent pattern matching",
            "mean(confidences) → overall confidence score",
        ],
        notes=[
            "3 anomaly patterns checked: supplier↔inventory cascade, shipment+supplier compound, inventory+shipment stockout",
            "anomaly_correlations exposed to frontend as distinct alert cards",
            "confidence = mean across agents — a single low-confidence agent drags the overall score down",
        ],
    )

    slide_node_detail(
        prs,
        icon="💡", title="Recommendation Engine Node",
        subtitle="Generates prioritised, domain-specific actionable recommendations",
        accent_col=GREEN,
        role_text=(
            "RecommendationAgent.analyze() receives the original query and all agent findings "
            "(including anomaly correlations from summary_node). Generates a list of actionable "
            "recommendations with priority, category, and rationale — drawn from all evidence "
            "across SQL data, vector-retrieved docs, and specialist agent analyses."
        ),
        inputs=[
            "query (str)",
            "agent_findings — all findings including anomalies",
            "execution_log",
        ],
        outputs=[
            "recommendations — list of {priority, category, action, rationale}",
            "execution_log — updated with count",
        ],
        calls=[
            "RecommendationAgent().analyze(query, findings) → LLM generation",
        ],
        notes=[
            "Recommendations are the final user-facing output alongside the answer",
            "Each recommendation tagged with: domain (supplier/shipment/inventory), priority (high/medium/low)",
            "Runs as last node before END in both paths: specialists ran OR SQL-only path",
            "Frontend RecommendationCard component renders each rec with colour-coded priority badges",
        ],
    )

    slide_anomaly_and_eval(prs)
    slide_tech_stack(prs)
    slide_business_value(prs)
    slide_demo(prs)
    slide_closing(prs)

    out_path = r"d:\FDE\Repo_Projects\capstone_ai_supply_chain_assistant\AI_Supply_Chain_Assistant_Presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
